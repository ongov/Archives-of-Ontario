"""
Set of helper functions that are needed in order to extract content from different formats of file.

@Author: Cheng Gao
@Date: Feb 5,2019
"""
import os # to execute command line operation
from gensim.summarization import keywords #for keyword extraction
from openpyxl import Workbook #to generate the result file
# the following import section is for docx files
import docx2txt
# the following section is for pdf files
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from io import BytesIO
# the following import section is for mp4 and wave files
import speech_recognition as sr
# the following import section is for pptx files
from pptx import Presentation
# the following import section is for object recognition
import keras
from keras.preprocessing.image import load_img, img_to_array
from keras.applications.imagenet_utils import decode_predictions
from keras.applications import nasnet
import numpy as np
import matplotlib.pyplot as plt
# the following import section is for pdf files that contains picture of texts
from wand.image import Image
from wand.color import Color
# the following import section is for message file
import extract_msg
# the following import section is for keyword refinement
from nltk.stem import PorterStemmer
from nltk.tokenize import word_tokenize
from Summarizer import TextSummarizer
# the following import section is for checksum
import hashlib
import textract

class infoGetter():
    " this class provides multiple methods to extract content from different file formats"
    def flattenPaths(self, rootDir):
        """
        assumptions: rootDir is a valid directory
        output: a tuple with the flattened paths
        """
        # create fileList
        filePaths =[]
        for dirName, subdirList, fileList in os.walk(rootDir):
            for fname in fileList:
                filePaths.append(dirName+'/'+fname)   
        return (filePaths)
    def png2txt(self, filePath):
        " this method performs optical character recognition"
        os.system("tesseract "+filePath.replace(' ','\ ')+ " file -l eng")
        f = open("file.txt", "rt", encoding = 'utf-8')
        content = f.read()
        os.remove("file.txt")
        return content
    def pdf2png2txt(self, filename, output_path, resolution=150):
        """ Convert a PDF into images.
    
            All the pages will give a single png file with format:
            {pdf_filename}-{page_number}.png
    
            The function removes the alpha channel from the image and
            replace it with a white background.
        """
        all_pages = Image(filename=filename, resolution=resolution)
        page_number = 0
        for i, page in enumerate(all_pages.sequence):
            page_number = i
            with Image(page) as img:
                img.format = 'png'
                img.background_color = Color('white')
                img.alpha_channel = 'remove'
    
                image_filename = os.path.splitext(os.path.basename(filename))[0]
                image_filename = '{}-{}.png'.format(image_filename, i)
                image_filename = os.path.join(output_path, image_filename)

                img.save(filename=image_filename)
        text = ""
        new_file_path = os.path.join(os.path.abspath(os.path.dirname(__file__)), 'uploads')
        for i in range(page_number):
            image_filename = os.path.splitext(os.path.basename(filename))[0]
            image_filename = '{}-{}.png'.format(image_filename, i)            
            text += self.png2txt(os.path.join(new_file_path, image_filename))
        return text    
    def jpg2txt(self, filePath):
        " this method detects the object in the image"
        # Load the NASNET Model
        nasnet_model = nasnet.NASNetMobile(weights='imagenet')        
        filename = filePath
        # Load an image in PIL format
        original = load_img(filename, target_size=(224, 224))
        print('PIL image size',original.size)
        plt.imshow(original)
        plt.show()
         
        # convert the PIL image to a numpy array
        # IN PIL - image is in (width, height, channel)
        # In Numpy - image is in (height, width, channel)
        numpy_image = img_to_array(original)
        plt.imshow(np.uint8(numpy_image))
        plt.show()
        print('numpy array size',numpy_image.shape)
         
        # Convert the image / images into batch format
        # expand_dims will add an extra dimension to the data at a particular axis
        # We want the input matrix to the network to be of the form (batchsize, height, width, channels)
        # Thus we add the extra dimension to the axis 0.
        image_batch = np.expand_dims(numpy_image, axis=0)
        print('image batch size', image_batch.shape)
        plt.imshow(np.uint8(image_batch[0]))       
        processed_image = nasnet.preprocess_input(image_batch.copy())
        predictions = nasnet_model.predict(processed_image)
        label = decode_predictions(predictions)
        text = "The object detected in this image is a "+ label[0][0][1]
        return text
    def pptx2txt(self, filePath):
        "this method extracts text content from powerpoints"
        text = ""
        prs = Presentation(filePath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr (shape, "text"):
                    text += shape.text + "\n"
        return text
    def mp42txt(self, filePath):
        "this method extracts text content from mp4 files"
        # convert the mp4 file to aac file first
        os.system("ffmpeg -i "+filePath.replace(" ","\ ")+ " -vn -acodec copy sound.aac")
        # convert the mp4 file to wave file
        os.system("ffmpeg -i sound.aac sound.wav")
        r = sr.Recognizer()
        with sr.AudioFile("sound.wav") as source:
            audio = r.record(source)
        text = r.recognize_sphinx(audio)
        os.remove("sound.aac")
        os.remove("sound.wav")
        return text
    def wav2txt(self, filePath):
        "this method extracts text content from wav files"
        r = sr.Recognizer()
        with sr.AudioFile(filePath) as source:
            audio = r.record(source)
        text = r.recognize_sphinx(audio)
        return text
    def doc2txt(self, filePath):
        "this method extracts text content from doc files"
        # convert doc file to txt file using antiword
        os.system("antiword "+ filePath.replace(" ","\ ") + " > " "file.txt")
        f = open("file.txt", "rt", errors= "replace")
        content = f.read()
        os.remove("file.txt")
        return content
    def pdf2txt(self, filePath):
        "this method extracts text content from pdf files"
        manager = PDFResourceManager()
        retstr = BytesIO()
        layout = LAParams(all_texts=True)
        device = TextConverter(manager, retstr, laparams=layout)
        filepath = open(filePath, 'rb')
        interpreter = PDFPageInterpreter(manager, device)
    
        for page in PDFPage.get_pages(filepath, check_extractable=True):
            interpreter.process_page(page)
    
        text = retstr.getvalue()
    
        filepath.close()
        device.close()
        retstr.close()
        return str(text)

    def xlsx2txt(self,filePath):
        return str(textract.process(filePath))
    def xls2txt(self,filePath):
        return str(textract.process(filePath))
    def txt2txt(self,filePath):
        file = open(filePath,'r')
        content = ""
        for line in file.readlines():
            content += line
        return content
    def getFileInfo(self, fileList):
        failedFile = []
        functionDict = {"doc":self.doc2txt, "docx":docx2txt.process,
                        "pdf":self.pdf2txt, "mp4": self.mp42txt,
                        "wav":self.wav2txt, "pptx":self.pptx2txt,
                        "jpg":self.jpg2txt, "png":self.png2txt,
                        "xlsx":self.xlsx2txt,"xls":self.xls2txt,
                        "txt":self.txt2txt}
        wb = Workbook()
        ws= wb.active
        ws['A1'] = 'File Name'
        ws['B1'] = 'File Title'
        ws['C1'] = 'Keywords'
        ws['D1'] = 'Summary'
        ws['E1'] = 'Checksum'
        ws['F1'] = 'Object Recognition'
        ws['G1'] = 'Optical Character Recognition'
        ws['H1'] = 'File Format'
        ws['I1'] = 'File Size'
        ws['J1'] = 'Status'
        ws['K1'] = 'Failed File'
        ws['L1'] = 'Reason'
        
        for filePath in fileList:
                text=''
                summary = ''
                checksum = ''
                try:
                    # call corresponding method based on the extension of the file
                    text = functionDict[os.path.splitext(filePath)[1][1:]](filePath)
                    if keywords(text,lemmatize=True).split('\n') == [''] and os.path.splitext(filePath)[1][1:] == 'pdf':
                        text = self.pdf2png2txt(filePath,os.path.join(os.path.abspath(os.path.dirname(__file__)),'uploads'),150)
                    text = text.replace('\\n', ' ')
                    checksum = hashlib.md5(text.encode('utf-8')).hexdigest()
                    print("Sucessfully extracted text from file: %s" % filePath)
                except Exception as e:
                    print(e)
                    failedFile.append(filePath)
                    print("Failed to read file: %s" % filePath)
                sentenceList=text.split('\n')
                print(sentenceList)
                
                # check file format
                fileFormat = os.path.splitext(filePath)[1][1:]
                fileName = os.path.splitext(filePath.split("/")[-1])[0]
                tags = ''
                # to find root word
                ps = PorterStemmer()                
                keyword = keywords(text,lemmatize=True).split('\n')
                status = 'successfully extracted keywords'
                if len(keyword)>5:
                    keyword = keyword[:5]
                if len(keyword) != 0:
                    for tag in keyword:
                        if keyword.index(tag) != len(keyword)-1:
                            tag = ps.stem(tag)
                            tags += tag + '; '
                        else:
                            tag = ps.stem(tag)
                            tags += tag
                if tags == '':
                    tags = 'N/A'
                    status = 'failed to extract keywords since there is no content'
                if tags != '':
                    # get summary
                    textSummarizer = TextSummarizer(text)
                    number = 1
                    for i in textSummarizer.getSummary(3):
                        summary += str(number) +': ' + i + ' '
                        number += 1
                if fileFormat != 'jpg' and fileFormat != 'png':
                    print(fileFormat)
                    if sentenceList[0] != '':
                        fileTitle = str(sentenceList[0])
                    else:
                        fileTitle = 'N/A'
                        status += ' failed to find title'
                    objRecog = 'N/A'
                    charRecog = 'N/A'
                else:
                    fileTitle = 'N/A'
                    status += ' failed to find title'
                    if fileFormat == 'jpg':
                        objRecog = text
                        charRecog = 'N/A'
                        tags = 'N/A' 
                    else:
                        charRecog = 'Here are the first 10 lines of the file:\n'
                        for line in sentenceList[:9]:
                            charRecog += str(line) + '\n'
                        objRecog = 'N/A'
                fileSize = os.path.getsize(filePath)/1000
                if fileSize > 100:
                    fileSize = str(round(fileSize/1000, 1)) + ' MB'
                else:
                    fileSize = str(round(fileSize, 1) ) + ' KB'
                
                try:
                    ws.append([fileName+'.'+fileFormat, fileTitle, tags, summary, checksum, objRecog, charRecog, fileFormat, fileSize, status])
                except Exception as e:
                    print(e)
                    print("Failed to get keywords: "+filePath)
        rowNum = 2
        for file in failedFile:
            ws['K'+str(rowNum)] = file
            ws['L'+str(rowNum)] = 'Failed to extract content from this file'
            rowNum+=1
        wb.save("result.xlsx")
